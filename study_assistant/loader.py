from io import BytesIO
from typing import BinaryIO, Optional

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def _read_pdf(file_like: BinaryIO) -> str:
	reader = PdfReader(file_like)
	texts = []
	for page in reader.pages:
		try:
			texts.append(page.extract_text() or "")
		except Exception:
			continue
	return "\n".join(texts)


def _read_docx(file_like: BinaryIO) -> str:
	doc = Document(file_like)
	return "\n".join(p.text for p in doc.paragraphs if p.text)


def _read_pptx(file_like: BinaryIO) -> str:
	prs = Presentation(file_like)
	texts = []
	for slide in prs.slides:
		for shape in slide.shapes:
			if hasattr(shape, "has_text_frame") and shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					text = "".join(run.text for run in paragraph.runs)
					if text:
						texts.append(text)
	return "\n".join(texts)


def _read_text(file_like: BinaryIO) -> str:
	data = file_like.read()
	if isinstance(data, bytes):
		return data.decode("utf-8", errors="ignore")
	return str(data)


def load_uploaded_file(uploaded_file: BinaryIO, filename: Optional[str]) -> str:
	"""Load text from a Streamlit UploadedFile or any file-like object given a filename."""
	name = (filename or "").lower()
	if name.endswith(".pdf"):
		return _read_pdf(uploaded_file)
	elif name.endswith(".docx"):
		return _read_docx(uploaded_file)
	elif name.endswith(".pptx"):
		return _read_pptx(uploaded_file)
	elif name.endswith(".txt") or name.endswith(".md"):
		return _read_text(uploaded_file)
	else:
		raise ValueError(f"Format non supporté: {filename}")


def load_path(path: str) -> str:
	"""Load text from a path on disk."""
	lower = path.lower()
	with open(path, "rb") as f:
		if lower.endswith(".pdf"):
			return _read_pdf(f)
		elif lower.endswith(".docx"):
			return _read_docx(f)
		elif lower.endswith(".pptx"):
			return _read_pptx(f)
		elif lower.endswith(".txt") or lower.endswith(".md"):
			return _read_text(f)
		else:
			raise ValueError(f"Format non supporté: {path}")